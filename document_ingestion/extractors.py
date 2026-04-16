"""Text extraction from various document formats (PDF, DOCX, PPTX)."""

import os
from pathlib import Path
from typing import List, Dict, Optional, Tuple

try:
    import fitz as pymupdf
except ImportError:
    try:
        import pymupdf
    except ImportError:
        pymupdf = None

try:
    from docx import Document
except ImportError:
    pass

try:
    from pptx import Presentation
except ImportError:
    pass


class PDFExtractor:
    """Extracts text and layout from PDFs using PyMuPDF (fitz)."""

    def __init__(self, min_chars_per_page: int = 50):
        self.min_chars_per_page = min_chars_per_page

    def extract(self, file_path: str) -> Tuple[List[Dict], bool]:
        """
        Extract text from PDF.
        Returns:
            Tuple of (list of page dicts, requires_ocr bool)
            Page dict: {"page_num": int, "text": str, "needs_ocr": bool}
        """
        if pymupdf is None:
            raise ImportError(
                "PyMuPDF is not installed. Run: pip install pymupdf"
            )

        doc = pymupdf.open(file_path)
        pages = []
        any_needs_ocr = False

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            
            # If page is practically empty of text, it might be a scanned image
            needs_ocr = len(text) < self.min_chars_per_page
            if needs_ocr:
                any_needs_ocr = True
                
            pages.append({
                "page_num": page_num + 1,
                "text": text,
                "needs_ocr": needs_ocr
            })

        doc.close()
        return pages, any_needs_ocr


class DOCXExtractor:
    """Extracts text from Word documents."""

    def extract(self, file_path: str) -> List[Dict]:
        """Returns list of chunk dicts (treating paragraphs as chunks)."""
        try:
            doc = Document(file_path)
        except NameError:
            raise ImportError("python-docx is not installed. Run: pip install python-docx")

        chunks = []
        current_text = []
        current_heading = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Very basic heading logic (can be refined via styles)
            if para.style.name.startswith('Heading'):
                # Save previous chunk
                if current_text:
                    chunks.append({
                        "section_heading": current_heading,
                        "text": "\n".join(current_text)
                    })
                    current_text = []
                current_heading = text
            else:
                current_text.append(text)

        if current_text:
            chunks.append({
                "section_heading": current_heading,
                "text": "\n".join(current_text)
            })

        return chunks


class PPTXExtractor:
    """Extracts text from PowerPoint presentations."""

    def extract(self, file_path: str) -> List[Dict]:
        """Returns list of page dicts (slides)."""
        try:
            prs = Presentation(file_path)
        except NameError:
            raise ImportError("python-pptx is not installed. Run: pip install python-pptx")

        pages = []

        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text.strip())
            
            text = "\n".join([t for t in text_runs if t])
            if text:
                pages.append({
                    "page_num": i + 1,
                    "text": text,
                    "needs_ocr": False # PPTX doesn't usually need full page OCR
                })

        return pages


def get_extractor(file_path: str):
    """Factory to get the right extractor based on extension."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return PDFExtractor()
    elif ext in [".docx", ".doc"]:
        return DOCXExtractor()
    elif ext in [".pptx", ".ppt"]:
        return PPTXExtractor()
    else:
        return None
